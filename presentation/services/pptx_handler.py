from pptx import Presentation


def create_presentation_pptx(presentation):
    prs = Presentation(presentation.template.file.path)
    for slide_data in presentation.slides.all():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.title
        slide.placeholders[1].text = slide_data.content

    pptx_path = f"media/{presentation.title}.pptx"
    prs.save(pptx_path)
    return pptx_path
